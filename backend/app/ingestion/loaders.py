"""Document loaders.

A loader's job is extraction plus position — never chunking. Each returns
`ParsedUnit`s carrying the text and enough metadata to render a citation a
human can verify: page, sheet and row range, slide, or heading path.

Units sharing a `group` may later be packed into one chunk. Units in different
groups never merge, because that would blur the citation.
"""

from __future__ import annotations

import csv
import io
import json
from collections.abc import Callable
from dataclasses import dataclass, field
from html.parser import HTMLParser
from typing import Any

from app.core.errors import ValidationFailed

# rows packed into a single tabular unit, header repeated in each
TABLE_ROWS_PER_UNIT = 40


@dataclass
class ParsedUnit:
    text: str
    group: str
    location: str
    metadata: dict[str, Any] = field(default_factory=dict)
    atomic: bool = False


@dataclass
class LoadResult:
    units: list[ParsedUnit]
    page_count: int | None = None
    source_type: str = "text"


# ------------------------------------------------------------------- helpers
def _clean(text: str) -> str:
    lines = [ln.rstrip() for ln in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")]
    out: list[str] = []
    blank = 0
    for ln in lines:
        if ln.strip():
            blank = 0
            out.append(ln)
        else:
            blank += 1
            if blank <= 1:
                out.append("")
    return "\n".join(out).strip()


def _rows_to_units(
    rows: list[list[str]], *, sheet: str, source_type: str
) -> list[ParsedUnit]:
    """Header row repeated in every unit — a table fragment without its header
    is unusable to the model and unverifiable to the reader."""
    if not rows:
        return []
    header = [str(c) if c is not None else "" for c in rows[0]]
    body = rows[1:]
    if not body:
        body = [header]
        header = [f"col_{i + 1}" for i in range(len(header))]

    units: list[ParsedUnit] = []
    for start in range(0, len(body), TABLE_ROWS_PER_UNIT):
        block = body[start : start + TABLE_ROWS_PER_UNIT]
        first, last = start + 2, start + len(block) + 1  # 1-based, header is row 1
        lines = [" | ".join(header)]
        for row in block:
            cells = ["" if c is None else str(c) for c in row]
            if not any(c.strip() for c in cells):
                continue
            lines.append(" | ".join(cells))
        if len(lines) == 1:
            continue
        location = f"Sheet '{sheet}' rows {first}-{last}" if sheet else f"Rows {first}-{last}"
        units.append(
            ParsedUnit(
                text="\n".join(lines),
                group=f"{sheet}:{start}",
                location=location,
                metadata={"sheet": sheet, "row_start": first, "row_end": last,
                          "source_type": source_type},
                atomic=True,
            )
        )
    return units


# ------------------------------------------------------------------- loaders
def load_pdf(data: bytes) -> LoadResult:
    from pypdf import PdfReader

    try:
        reader = PdfReader(io.BytesIO(data))
    except Exception as exc:
        raise ValidationFailed(f"Could not read this PDF: {exc}") from exc

    if reader.is_encrypted:
        try:
            reader.decrypt("")
        except Exception:
            raise ValidationFailed(
                "This PDF is password protected. Upload an unlocked copy."
            ) from None

    units: list[ParsedUnit] = []
    for index, page in enumerate(reader.pages, start=1):
        try:
            text = _clean(page.extract_text() or "")
        except Exception:
            text = ""
        if not text:
            continue
        units.append(
            ParsedUnit(
                text=text,
                group=f"page:{index}",
                location=f"Page {index}",
                metadata={"page": index, "source_type": "pdf"},
            )
        )
    return LoadResult(units=units, page_count=len(reader.pages), source_type="pdf")


def load_docx(data: bytes) -> LoadResult:
    import docx

    document = docx.Document(io.BytesIO(data))
    units: list[ParsedUnit] = []
    heading_path: list[str] = []
    buffer: list[str] = []
    section_no = 0

    def flush() -> None:
        nonlocal buffer, section_no
        text = _clean("\n".join(buffer))
        buffer = []
        if not text:
            return
        section_no += 1
        path = " > ".join(heading_path) if heading_path else "Document body"
        units.append(
            ParsedUnit(
                text=text,
                group=f"section:{section_no}",
                location=path,
                metadata={"heading_path": list(heading_path), "source_type": "docx"},
            )
        )

    for para in document.paragraphs:
        style = (para.style.name or "") if para.style else ""
        text = para.text.strip()
        if style.startswith("Heading"):
            flush()
            try:
                level = int(style.split()[-1])
            except (ValueError, IndexError):
                level = 1
            heading_path[:] = heading_path[: level - 1]
            if text:
                heading_path.append(text)
            continue
        if text:
            buffer.append(text)
    flush()

    for t_index, table in enumerate(document.tables, start=1):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        units.extend(_rows_to_units(rows, sheet=f"Table {t_index}", source_type="docx"))

    return LoadResult(units=units, source_type="docx")


def load_xlsx(data: bytes) -> LoadResult:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    units: list[ParsedUnit] = []
    for sheet in workbook.worksheets:
        rows = [
            ["" if cell is None else str(cell) for cell in row]
            for row in sheet.iter_rows(values_only=True)
        ]
        rows = [r for r in rows if any(c.strip() for c in r)]
        units.extend(_rows_to_units(rows, sheet=sheet.title, source_type="xlsx"))
    workbook.close()
    return LoadResult(units=units, source_type="xlsx")


def load_csv(data: bytes) -> LoadResult:
    text = data.decode("utf-8", errors="replace")
    sample = text[:8192]
    try:
        dialect: Any = csv.Sniffer().sniff(sample, delimiters=",;\t|")
    except csv.Error:
        dialect = csv.excel
    rows = [row for row in csv.reader(io.StringIO(text), dialect) if any(c.strip() for c in row)]
    return LoadResult(units=_rows_to_units(rows, sheet="", source_type="csv"), source_type="csv")


def load_pptx(data: bytes) -> LoadResult:
    from pptx import Presentation

    deck = Presentation(io.BytesIO(data))
    units: list[ParsedUnit] = []
    for index, slide in enumerate(deck.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            frame = getattr(shape, "text_frame", None)
            if frame is not None:
                value = frame.text.strip()
                if value:
                    parts.append(value)
        text = _clean("\n".join(parts))
        if not text:
            continue
        units.append(
            ParsedUnit(
                text=text,
                group=f"slide:{index}",
                location=f"Slide {index}",
                metadata={"slide": index, "source_type": "pptx"},
            )
        )
    return LoadResult(units=units, page_count=len(deck.slides), source_type="pptx")


def load_json(data: bytes) -> LoadResult:
    try:
        parsed = json.loads(data.decode("utf-8", errors="replace"))
    except json.JSONDecodeError as exc:
        raise ValidationFailed(f"Invalid JSON: {exc.msg} (line {exc.lineno})") from exc

    def flatten(node: Any, path: str, out: list[str]) -> None:
        if isinstance(node, dict):
            for key, value in node.items():
                flatten(value, f"{path}.{key}" if path else str(key), out)
        elif isinstance(node, list):
            for i, value in enumerate(node):
                flatten(value, f"{path}[{i}]", out)
        else:
            out.append(f"{path}: {node}")

    units: list[ParsedUnit] = []
    top_level = parsed if isinstance(parsed, dict) else {"root": parsed}
    for key, value in top_level.items():
        lines: list[str] = []
        flatten(value, str(key), lines)
        text = _clean("\n".join(lines))
        if not text:
            continue
        units.append(
            ParsedUnit(
                text=text,
                group=f"key:{key}",
                location=f"Key '{key}'",
                metadata={"json_key": key, "source_type": "json"},
            )
        )
    return LoadResult(units=units, source_type="json")


class _TextExtractor(HTMLParser):
    _SKIP = {"script", "style", "noscript", "head"}

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self._skipping = 0

    def handle_starttag(self, tag: str, attrs: Any) -> None:  # noqa: ARG002
        if tag in self._SKIP:
            self._skipping += 1

    def handle_endtag(self, tag: str) -> None:
        if tag in self._SKIP and self._skipping:
            self._skipping -= 1
        elif tag in {"p", "div", "br", "li", "tr", "h1", "h2", "h3", "h4"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._skipping and data.strip():
            self.parts.append(data.strip())


def load_html(data: bytes) -> LoadResult:
    parser = _TextExtractor()
    parser.feed(data.decode("utf-8", errors="replace"))
    text = _clean(" ".join(parser.parts))
    units = (
        [ParsedUnit(text=text, group="body", location="Document",
                    metadata={"source_type": "html"})]
        if text
        else []
    )
    return LoadResult(units=units, source_type="html")


def load_text(data: bytes) -> LoadResult:
    text = _clean(data.decode("utf-8", errors="replace"))
    units = (
        [ParsedUnit(text=text, group="body", location="Document",
                    metadata={"source_type": "text"})]
        if text
        else []
    )
    return LoadResult(units=units, source_type="text")


# ------------------------------------------------------------------ registry
LOADERS: dict[str, Callable[[bytes], LoadResult]] = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".xlsx": load_xlsx,
    ".xlsm": load_xlsx,
    ".csv": load_csv,
    ".pptx": load_pptx,
    ".json": load_json,
    ".html": load_html,
    ".htm": load_html,
    ".txt": load_text,
    ".md": load_text,
    ".log": load_text,
}

MIME_TYPES: dict[str, str] = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".xlsm": "application/vnd.ms-excel.sheet.macroEnabled.12",
    ".csv": "text/csv",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".json": "application/json",
    ".html": "text/html",
    ".htm": "text/html",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".log": "text/plain",
}

SUPPORTED_EXTENSIONS = sorted(LOADERS)


def extension_of(filename: str) -> str:
    _, _, ext = filename.rpartition(".")
    return f".{ext.lower()}" if ext else ""


def load(filename: str, data: bytes) -> LoadResult:
    ext = extension_of(filename)
    loader = LOADERS.get(ext)
    if loader is None:
        raise ValidationFailed(
            f"Unsupported file type '{ext or filename}'. "
            f"Supported: {', '.join(SUPPORTED_EXTENSIONS)}"
        )
    return loader(data)
